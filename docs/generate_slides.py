#!/usr/bin/env python3
"""Generate Agentic AI SDLC prediction-engine deck (Google Slides–compatible .pptx)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


# --- Theme --------------------------------------------------------------------
NAVY_BG = RGBColor(13, 17, 23)
HEADER_ROW_BG = RGBColor(30, 41, 59)
TEXT_LIGHT = RGBColor(242, 242, 242)
ACCENT = RGBColor(56, 189, 248)
TEXT_MUTED = RGBColor(148, 163, 184)
FONT_NAME = "Calibri"

OUT_PATH = Path(__file__).resolve().parent / "agentic-ai-sdlc-prediction-engine.pptx"


def _set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BG


def _set_cell_border(cell, color: str = "64748b", width: str = "9525") -> None:
    """Light border on all sides (hex RGB without #)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("lnL", "lnT", "lnR", "lnB"):
        el = OxmlElement(f"a:{edge}")
        el.set("w", width)
        solid = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", color)
        solid.append(srgb)
        el.append(solid)
        tcPr.append(el)


def _style_run(run, *, size_pt: int, bold: bool = False, color: RGBColor = TEXT_LIGHT) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    *,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = auto_size
    tf.vertical_anchor = valign
    tf.paragraphs[0].alignment = align
    return tf


def _title_bar(slide, title: str) -> None:
    left, top, w = Inches(0.55), Inches(0.45), Inches(9.0)
    tf = _add_textbox(slide, left, top, w, Inches(0.85), align=PP_ALIGN.LEFT)
    p = tf.paragraphs[0]
    p.text = title
    _style_run(p.runs[0], size_pt=26, bold=True, color=ACCENT)


def _body_frame(slide, top, height=Inches(5.9)):
    left, w = Inches(0.7), Inches(8.85)
    return _add_textbox(slide, left, top, w, height, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP)


def add_bullets(slide, title: str, bullets: list[str]) -> None:
    _title_bar(slide, title)
    tf = _body_frame(slide, Inches(1.08))
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = FONT_NAME
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_LIGHT
        p.line_spacing = 1.12
        try:
            _style_run(p.runs[0], size_pt=16, color=TEXT_LIGHT)
        except IndexError:
            pass


def add_numbered(slide, title: str, lines: list[str]) -> None:
    _title_bar(slide, title)
    tf = _body_frame(slide, Inches(1.08))
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}. {line}"
        p.level = 0
        p.space_after = Pt(5)
        try:
            _style_run(p.runs[0], size_pt=15, color=TEXT_LIGHT)
        except IndexError:
            pass


def add_table_slide(slide, title: str, headers: list[str], rows: list[list[str]]) -> None:
    _title_bar(slide, title)
    rows_n, cols_n = len(rows) + 1, len(headers)
    left, top = Inches(0.65), Inches(1.05)
    width, height = Inches(9.1), Inches(0.42 * (rows_n + 1) + 0.15)
    table_shape = slide.shapes.add_table(rows_n, cols_n, left, top, width, height)
    tbl = table_shape.table

    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_ROW_BG
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            for r in p.runs:
                _style_run(r, size_pt=13, bold=True, color=ACCENT)
        _set_cell_border(cell)

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = text
            tf = cell.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf.paragraphs:
                for rr in p.runs:
                    _style_run(rr, size_pt=12, color=TEXT_LIGHT)
            _set_cell_border(cell)


def add_diagram_text(slide, title: str, diagram: str) -> None:
    _title_bar(slide, title)
    tf = _body_frame(slide, Inches(1.08), height=Inches(5.6))
    tf.clear()
    p = tf.paragraphs[0]
    p.text = diagram.strip("\n")
    for r in p.runs:
        _style_run(r, size_pt=12, color=TEXT_LIGHT)
        r.font.name = FONT_NAME


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    tf = _add_textbox(
        slide, Inches(0.9), Inches(2.15), Inches(8.2), Inches(1.15), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE
    )
    p = tf.paragraphs[0]
    p.text = "Agentic AI SDLC"
    _style_run(p.runs[0], size_pt=44, bold=True, color=TEXT_LIGHT)

    st = _add_textbox(
        slide, Inches(0.9), Inches(3.35), Inches(8.2), Inches(1.0), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP
    )
    p2 = st.paragraphs[0]
    p2.text = "Building the Right-Sizing Prediction Engine for ACM Observability"
    _style_run(p2.runs[0], size_pt=20, bold=False, color=ACCENT)

    ft = _add_textbox(
        slide, Inches(0.75), Inches(6.55), Inches(8.5), Inches(0.55), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM
    )
    pf = ft.paragraphs[0]
    pf.text = "ACM Observability Team | May 2026"
    _style_run(pf.runs[0], size_pt=14, color=TEXT_MUTED)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "The Challenge",
        [
            "ACM clusters need proactive resource management",
            "Reactive right-sizing misses trends and anomalies",
            "Need: Predict future resource needs across Namespace, Workload, GPU, VM levels",
            "Constraint: Zero data exfiltration — all processing on-cluster",
        ],
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "What is Agentic AI SDLC?",
        [
            "AI operates within constrained phases",
            "Each phase has explicit inputs, outputs, human checkpoints",
            "Core principle: Structure In, Structure Out",
            "Always grounded in real code (file paths, symbols, patterns)",
            'Never "implement this" from vague description',
        ],
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_numbered(
        slide,
        "The 8-Phase Workflow",
        [
            "Repository Impact Map — Human Checkpoint",
            "Structured Task Creation — Human Checkpoint",
            "Implementation (Parallel Execution)",
            "Testing",
            "PR Creation",
            "CI & Auto-fix",
            "Integration Testing",
            "Verification & Review",
        ],
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "Phase 1–2: Planning (Grounded in Real Code)",
        [
            "Scanned MCO + MCOA codebases before any code written",
            "Identified 50+ files, exact APIs, modules affected",
            "Structured tasks with real file paths, real symbol names",
            "Key decisions (human-approved):",
            "  • Prediction as sibling field under PlatformAnalyticsSpec",
            "  • Pluggable provider architecture (4 types)",
            "  • Multi-dimension support",
        ],
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "Phase 3: Parallel Implementation",
        [
            "MCO and MCOA implemented simultaneously via AI subagents",
            "MCO: API types, ADC sync, CRD bundle, scrape configs, dashboards",
            "MCOA: Multi-model ensemble, training controller, providers, Helm charts",
            "Multi-dimension: Namespace, Workload/Pod, GPU, Virtualization/VM",
        ],
    )

    # Slide 7
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    diagram = """
MCO CR (Hub) → MCO Operator → AddOnDeploymentConfig
    → MCOA Addon Manager (Hub)
        → Training Controller (Holt-Winters + STL + AR)
        → ManifestWork per spoke
    → Spoke: Config, NetworkPolicy, RBAC, Scrape Config
    ← Thanos (30d history, 6h training cycle)
""".strip("\n")
    add_diagram_text(slide, "Prediction Engine Architecture", diagram)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "Multi-Model Ensemble",
        [
            "Holt-Winters: Seasonal decomposition with exponential smoothing",
            "STL: Seasonal-Trend decomposition using Loess",
            "AR(p): Autoregressive model for short-term patterns",
            "Anomaly detection via ensemble agreement scoring",
            "Optimization recommender for resource right-sizing",
        ],
    )

    # Slide 9
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_table_slide(
        slide,
        "Pluggable Provider Architecture",
        ["Provider", "Description", "Data Privacy", "Dependencies"],
        [
            ["Built-in", "Go-native ensemble", "On-cluster", "None"],
            ["ONNX", "Pre-trained models", "On-cluster", "ONNX runtime"],
            ["External API", "Prediction service", "Requires consent", "API endpoint"],
            ["Custom", "User endpoint", "Requires consent", "Custom endpoint"],
        ],
    )

    # Slide 10
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_table_slide(
        slide,
        "4 Prediction Dimensions",
        ["Dimension", "Metrics", "Forecast Series"],
        [
            ["Namespace", "CPU, Memory", "acm_rs:prediction_forecast_cpu/memory"],
            ["Workload/Pod", "CPU, Memory", "acm_rs:prediction_forecast_workload_cpu/memory"],
            ["GPU", "Utilization, Memory", "acm_rs:prediction_forecast_gpu_utilization/memory"],
            ["VM", "CPU, Memory", "acm_rs:prediction_forecast_vm_cpu/memory"],
        ],
    )

    # Slide 11
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "Phase 6: CI & Auto-fix",
        [
            "96 golangci-lint issues fixed automatically (MCOA)",
            "45 missing scrape config metrics added (MCO)",
            "CRD bundle drift resolved (MCO)",
            "dashcheck prediction metrics exclusion (MCO)",
            "Total: 99 CI issues fixed across both repos",
        ],
    )

    # Slide 12
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    _title_bar(slide, "Phase 7: Integration Testing")
    tf = _body_frame(slide, Inches(1.08))
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "7 deployment blockers solved:"
    _style_run(p0.runs[0], size_pt=16, bold=True, color=ACCENT)
    blockers = [
        "mch-image-manifest reverted → MCO CR annotations",
        "Binary name mismatch → os.Args detection",
        "Read-only /tmp → emptyDir volume",
        "MCH reverts operator → image-override workflow",
        'ADC value mismatch → "enabled" not "true"',
        "Training not wired → hub_run.go",
        "Architecture confusion → MCOA runs on hub",
    ]
    for i, line in enumerate(blockers, 1):
        p = tf.add_paragraph()
        num_run = p.add_run()
        num_run.text = f"{i}. "
        _style_run(num_run, size_pt=14, bold=True, color=ACCENT)
        txt_run = p.add_run()
        txt_run.text = line
        _style_run(txt_run, size_pt=14, color=TEXT_LIGHT)
        p.level = 0
        p.space_after = Pt(3)

    # Slide 13
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "Data Privacy by Design",
        [
            "NetworkPolicy on every spoke (rs-prediction-policy)",
            "Built-in provider: zero external calls",
            "Training data: Thanos in-cluster queries only",
            "RBAC-scoped ServiceAccount",
            "dataExfiltrationConsent required for external providers",
            "100% on-cluster processing by default",
        ],
    )

    # Slide 14
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_table_slide(
        slide,
        "Results & Metrics",
        ["Metric", "Value"],
        [
            ["Phases completed", "8"],
            ["Repos modified", "2 (MCO + MCOA)"],
            ["CI issues fixed", "99"],
            ["Deploy blockers solved", "7"],
            ["Prediction dimensions", "4"],
            ["Provider types", "4"],
            ["Data on-cluster", "100%"],
            ["Training cycle", "6h interval, 30d history"],
        ],
    )

    # Slide 15
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "The Agentic AI Advantage",
        [
            "Parallel execution: MCO + MCOA simultaneously",
            "Real-time debugging: Complex cluster issues diagnosed iteratively",
            "Code quality: Lint, tests, CRD validation automated",
            "Human–AI collaboration: Humans decide architecture, AI executes",
            "Phase discipline: No scope creep, no architectural freelancing",
        ],
    )

    # Slide 16
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "Lessons Learned",
        [
            "Always ground plans in real code before implementation",
            "Understand full deployment chain (MCH → MCO → MCOA → ManifestWork)",
            "Phase discipline prevents scope creep",
            "Human checkpoints catch architectural issues early",
            "Parallel AI agents reduce wall-clock time significantly",
        ],
    )

    # Slide 17
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    add_bullets(
        slide,
        "What's Next",
        [
            "Monitor training controller for first predictions",
            "Verify Perses dashboards with real data",
            "Open upstream PRs when ready",
            "E2E test suite for prediction pipeline",
            "Expand to additional prediction models",
        ],
    )

    # Slide 18
    slide = prs.slides.add_slide(blank)
    _set_slide_background(slide)
    tf = _add_textbox(
        slide, Inches(0.9), Inches(2.65), Inches(8.2), Inches(1.1), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE
    )
    p = tf.paragraphs[0]
    p.text = "Questions?"
    _style_run(p.runs[0], size_pt=40, bold=True, color=TEXT_LIGHT)
    st = _add_textbox(
        slide, Inches(0.9), Inches(3.85), Inches(8.2), Inches(0.95), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP
    )
    p2 = st.paragraphs[0]
    p2.text = "Agentic AI SDLC — Structure In, Structure Out"
    _style_run(p2.runs[0], size_pt=18, bold=False, color=ACCENT)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()